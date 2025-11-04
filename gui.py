from pipeline import *  # Import des fonctions corrigées
from pathlib import Path
import numpy as np
import cv2, os

# Import necessary PyQt5 components
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QLabel, 
    QCheckBox, QGroupBox, QRadioButton, QLineEdit, QSlider, QProgressBar, 
    QComboBox, QTabWidget, QPlainTextEdit, QTableWidget, QSplitter, QAction, 
    QToolBar, QFileDialog, QMessageBox, QHeaderView, QTableWidgetItem, QPushButton,
    QButtonGroup
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, pyqtSlot
from PyQt5.QtGui import QPixmap, QImage
from skimage.filters import frangi, meijering
from skimage.morphology import closing, remove_small_objects, square

# Import plotting libraries
import matplotlib.pyplot as plt
from matplotlib.backends.backend_qt5agg import FigureCanvasQTAgg as FigureCanvas

# Import utility libraries
import sys
import csv
from io import BytesIO

# Conditional imports for export functionalities
try:
    from reportlab.pdfgen import canvas as pdf_canvas
    from reportlab.lib.pagesizes import letter
    PDF_OK = True
except ImportError:
    PDF_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from prettytable import PrettyTable
    PRETTYTABLE_OK = True
except ImportError:
    PRETTYTABLE_OK = False

def numpy_to_qpixmap(array):
    """Convert numpy array (BGR or Grayscale) to QPixmap."""
    if array is None: return QPixmap()
    if len(array.shape) == 3: # BGR
        h, w, ch = array.shape
        fmt = QImage.Format_BGR888
    elif len(array.shape) == 2: # Grayscale
        h, w = array.shape
        ch = 1
        fmt = QImage.Format_Grayscale8
        array = np.require(array, np.uint8, 'C') # Ensure correct format
    else:
        return QPixmap() # Unsupported format
        
    bytes_per_line = ch * w
    qimg = QImage(array.data, w, h, bytes_per_line, fmt)
    return QPixmap.fromImage(qimg)

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Retina-ID | LBP + LDA (Corrigé)") # Titre mis à jour
        self.resize(1380, 900)
        self.opts = (True,True,True,True)
        self.g_rep = {} # Stores evaluation results (populated ONLY after test evaluation)
        self.models = {} # Stores trained models
        self.curr_model_name = "LDA"
        self.threshold = 0.8  # Default threshold
        self.train_folder = None
        self.trained_labels = [] # Store labels from training set

        self._ui()
        self.setStyleSheet(self._stylesheet())
        self._clear_evaluation_display() # Clear metrics display initially

    def _ui(self):
        self._actions()
        self._toolbar()
        central = QWidget(); layC = QHBoxLayout(central); self.setCentralWidget(central)
        splitter = QSplitter(Qt.Horizontal); layC.addWidget(splitter)

        # --- Left Panel --- 
        left = QWidget(); vL = QVBoxLayout(left); splitter.addWidget(left)
        self.lbl_orig = QLabel("Image originale"); self.lbl_orig.setFixedSize(320,320)
        self.lbl_orig.setAlignment(Qt.AlignCenter); self._frame(self.lbl_orig,2)
        vL.addWidget(self.lbl_orig, alignment=Qt.AlignCenter)

        opts_layout = QHBoxLayout()
        self.chk_c = QCheckBox("CLAHE");   self.chk_c.setChecked(True)
        self.chk_d = QCheckBox("Denoise"); self.chk_d.setChecked(True)
        self.chk_v = QCheckBox("Vaisseaux"); self.chk_v.setChecked(True)
        self.chk_m = QCheckBox("Morpho");  self.chk_m.setChecked(True)
        for w in (self.chk_c,self.chk_d,self.chk_v,self.chk_m): opts_layout.addWidget(w)
        vL.addLayout(opts_layout)

        grp_steps = QGroupBox("Étapes"); h_steps = QHBoxLayout(grp_steps)
        self.lbl_bin = self._step_label(); self.lbl_skel=self._step_label(); self.lbl_lbp=self._step_label()
        for l_step in (self.lbl_bin,self.lbl_skel,self.lbl_lbp): h_steps.addWidget(l_step)
        vL.addWidget(grp_steps)

        grpM = QGroupBox("Mode"); vM = QVBoxLayout(grpM)
        self.rb_ident = QRadioButton("Identification"); self.rb_ident.setChecked(True)
        self.rb_auth  = QRadioButton("Authentification")
        self.line_claim = QLineEdit(); self.line_claim.setPlaceholderText("Identité revendiquée")
        self.line_claim.setEnabled(False)
        self.rb_auth.toggled.connect(lambda v: self.line_claim.setEnabled(v))
        self.mode_group = QButtonGroup(); self.mode_group.addButton(self.rb_ident); self.mode_group.addButton(self.rb_auth)
        vM.addWidget(self.rb_ident); vM.addWidget(self.rb_auth); vM.addWidget(self.line_claim)
        
        threshold_layout = QHBoxLayout()
        threshold_layout.addWidget(QLabel("Seuil:"))
        self.threshold_slider = QSlider(Qt.Horizontal)
        self.threshold_slider.setMinimum(0)
        self.threshold_slider.setMaximum(100)
        self.threshold_slider.setValue(int(self.threshold * 100))
        self.threshold_slider.setTickPosition(QSlider.TicksBelow)
        self.threshold_slider.setTickInterval(10)
        self.threshold_slider.valueChanged.connect(self._update_threshold)
        self.threshold_label = QLabel(f"{self.threshold:.2f}")
        threshold_layout.addWidget(self.threshold_slider)
        threshold_layout.addWidget(self.threshold_label)
        vM.addLayout(threshold_layout)
        vL.addWidget(grpM)

        # --- Right Panel --- 
        right = QWidget(); vR = QVBoxLayout(right); splitter.addWidget(right)
        self.pbar = QProgressBar(); vR.addWidget(self.pbar)
        
        dataset_buttons_layout = QHBoxLayout()
        self.combo_method = QComboBox()
        self.combo_method.addItems(["LDA"])
        self.combo_method.currentTextChanged.connect(self._switch_method)
        dataset_buttons_layout.addWidget(self.combo_method, alignment=Qt.AlignLeft)
        
        self.btn_evaluate_test = QPushButton("Évaluer sur test (TFA/TFR/TEE)") # Label mis à jour
        self.btn_evaluate_test.setEnabled(False) # Disabled until model is trained
        self.btn_evaluate_test.clicked.connect(self._evaluate_test_folder)
        dataset_buttons_layout.addWidget(self.btn_evaluate_test)
        vR.addLayout(dataset_buttons_layout)

        self.tabs = QTabWidget(); vR.addWidget(self.tabs)
        self.txt_logs = QPlainTextEdit(readOnly=True)
        self.tabs.addTab(self.txt_logs, "Logs")

        tab_tbl = QWidget(); vTbl = QVBoxLayout(tab_tbl)
        self.tbl_metrics = QTableWidget(); self.tbl_cm = QTableWidget()
        vTbl.addWidget(QLabel("Métriques biométriques par classe (sur Test)")); vTbl.addWidget(self.tbl_metrics)
        vTbl.addWidget(QLabel("Matrice de confusion (sur Test)")); vTbl.addWidget(self.tbl_cm)
        self.tabs.addTab(tab_tbl, "Tableaux")

        tab_graph = QWidget(); vG = QVBoxLayout(tab_graph)
       
        self.canvas_det = FigureCanvas(plt.Figure(figsize=(4,4)))
        self.canvas_hist = FigureCanvas(plt.Figure(figsize=(4,3)))
        self.canvas_tfa = FigureCanvas(plt.Figure(figsize=(4,4)))
        self.canvas_tfr = FigureCanvas(plt.Figure(figsize=(4,4)))
        self.canvas_tfa = FigureCanvas(plt.Figure(figsize=(4, 4)))
        self.canvas_tfr = FigureCanvas(plt.Figure(figsize=(4, 4)))
        vG.addWidget(QLabel("Courbe DET (TFA/TFR/TEE sur Test)")); vG.addWidget(self.canvas_det)
        vG.addWidget(QLabel("Courbe TFA (Test)"))
        vG.addWidget(self.canvas_tfa)

        vG.addWidget(QLabel("Courbe TFR (Test)"))
        vG.addWidget(self.canvas_tfr) 


       
        self.tabs.addTab(tab_graph, "Graphes")

        tab_ov = QWidget(); vO = QVBoxLayout(tab_ov)
        self.txt_stats = QPlainTextEdit(readOnly=True)
        self.tbl_feats = QTableWidget()
        vO.addWidget(QLabel("Statistiques globales (sur Test)")); vO.addWidget(self.txt_stats)
        
     
        self.tabs.addTab(tab_ov, "Aperçu modèle")

        splitter.setSizes([640,740])

    def _validate_retina_image_advanced(
        self,
        img: np.ndarray,
        min_ratio: float = 0.05,
        downscale: float = 0.5,
        vessel_thresh: float = 0.15,
        frangi_scales: tuple = (1, 8),
        frangi_step: float = 2.0,
        clahe_clip: float = 0.03,
        min_obj_size: int = 64
    ) -> bool:
        """
        Validation “retine‐like” optimisée (Frangi + Meijering + CLAHE + morpho).
        """
        # 1) Sous‐échantillonnage (optionnel)
        if downscale != 1.0:
            img = cv2.resize(img, None,
                             fx=downscale, fy=downscale,
                             interpolation=cv2.INTER_AREA)

        # 2) Gris + CLAHE
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY) if img.ndim == 3 else img
        clahe = cv2.createCLAHE(clipLimit=clahe_clip, tileGridSize=(8, 8))
        gray_eq = clahe.apply(gray)

        # 3) Vesselness multi‐échelle
        v1 = frangi(gray_eq,
                    scale_range=frangi_scales,
                    scale_step=frangi_step,
                    black_ridges=False)
        v2 = meijering(gray_eq,
                        sigmas=np.arange(frangi_scales[0],
                                         frangi_scales[1]+1,
                                         frangi_step),
                        black_ridges=False)
        vesselness = np.maximum(v1, v2)

        # 4) Seuillage + nettoyage morpho
        mask = vesselness > vessel_thresh
        mask = closing(mask, square(3))
        mask = remove_small_objects(mask, min_size=min_obj_size)

        # 5) Calcul du ratio
        ratio = mask.sum() / mask.size

        return ratio >= min_ratio    

    def _update_threshold(self, value):
        self.threshold = value / 100.0
        self.threshold_label.setText(f"{self.threshold:.2f}")
        # If results exist, re-evaluate with new threshold
        if self.g_rep and self.curr_model_name in self.g_rep and self.g_rep[self.curr_model_name]:
             self._evaluate_test_folder() # Re-run evaluation with new threshold

    def _actions(self):
        self.act_open_ds  = QAction("Charger dataset & Entraîner", self, triggered=self._load_dataset)
        self.act_open_img = QAction("Charger image seule...", self, triggered=self._prompt_open_image)
        self.act_analyse  = QAction("Analyser image", self, triggered=self._analyse)
        self.act_analyse.setEnabled(False) # Disabled until model is trained
        self.act_export_csv = QAction("Export CSV (Test)", self, triggered=self._export_csv)
        self.act_export_pdf = QAction("Export PDF (Test)", self, triggered=self._export_pdf)
        self.act_export_ppt = QAction("Export PPTX (Test)", self, triggered=self._export_ppt)
        for act in [self.act_export_csv, self.act_export_pdf, self.act_export_ppt]:
            act.setEnabled(False) # Disabled until evaluation is done

    def _toolbar(self):
        tb = QToolBar(); self.addToolBar(tb)
        for a in (self.act_open_ds, self.act_open_img, self.act_analyse,
                  self.act_export_csv, self.act_export_pdf, self.act_export_ppt):
            tb.addAction(a)

    def _frame(self, lbl, w=1):
        lbl.setStyleSheet(f"border:{w}px solid #4caf50; border-radius:6px;")

    def _step_label(self):
        l = QLabel(); l.setFixedSize(200,200); l.setAlignment(Qt.AlignCenter)
        self._frame(l,1); return l

    def _load_dataset(self):
        folder = QFileDialog.getExistingDirectory(self,"Sélectionner le dossier RIBD (contenant train/ et test/)","./RIBD")
        if not folder: return
        
        train_folder_path = os.path.join(folder, "train")
        test_folder_path = os.path.join(folder, "test")
        
        if not os.path.exists(train_folder_path) or not os.path.isdir(train_folder_path):
            QMessageBox.critical(self, "Erreur", f"Dossier 'train' non trouvé dans {folder}")
            return
        if not os.path.exists(test_folder_path) or not os.path.isdir(test_folder_path):
            QMessageBox.critical(self, "Erreur", f"Dossier 'test' non trouvé dans {folder}")
            return
            
        self.train_folder = folder # Store base folder
        self.txt_logs.appendPlainText(f"→ Chargement du dossier train ({train_folder_path}) et entraînement...")
        self.pbar.setValue(0)
        self._clear_evaluation_display() # Clear previous results
        self.btn_evaluate_test.setEnabled(False) # Disable eval button during training
        self.act_analyse.setEnabled(False)
        for act in [self.act_export_csv, self.act_export_pdf, self.act_export_ppt]:
            act.setEnabled(False)
            
        opts = (self.chk_c.isChecked(), self.chk_d.isChecked(),
                self.chk_v.isChecked(), self.chk_m.isChecked())
                
        # Utilisation du LoaderThread de pipeline_corrige
        self.loader = LoaderThread(train_folder_path, opts, self)
        self.loader.progress.connect(self.pbar.setValue)
        self.loader.finished.connect(self._training_ready)
        self.loader.start()

    @pyqtSlot(object, object) # Modifié pour accepter g_rep et models
    def _training_ready(self, g_rep_train, trained_models):
        self.pbar.setValue(100)
        if "error" in g_rep_train:
            QMessageBox.critical(self,"Erreur d'entraînement", g_rep_train["error"]); return
            
        self.models = trained_models
        # Les labels sont maintenant dans g_rep_train['LDA']['labels']
        if self.curr_model_name in g_rep_train and 'labels' in g_rep_train[self.curr_model_name]:
            self.trained_labels = g_rep_train[self.curr_model_name]['labels']
        else:
             QMessageBox.warning(self, "Erreur", "Labels non trouvés après l'entraînement.")
             self.trained_labels = []
             
        self.txt_logs.appendPlainText("✓ Modèle(s) entraîné(s) sur le dossier train.")
        self.txt_logs.appendPlainText("→ Prêt à évaluer sur le dossier test.")
        self.btn_evaluate_test.setEnabled(True) # Enable evaluation button
        self.act_analyse.setEnabled(True) # Enable single image analysis
        self._update_lda_overview() # Update LDA overview (coefficients) as it depends only on training

    def _evaluate_test_folder(self):
        if not self.train_folder:
            QMessageBox.warning(self, "Erreur", "Veuillez d'abord charger un dataset et entraîner le modèle.")
            return
        if not self.models:
            QMessageBox.warning(self, "Erreur", "Aucun modèle entraîné disponible.")
            return
            
        test_folder_path = os.path.join(self.train_folder, "test")
        if not os.path.exists(test_folder_path) or not os.path.isdir(test_folder_path):
            QMessageBox.critical(self, "Erreur", f"Dossier test non trouvé: {test_folder_path}")
            return
            
        self.txt_logs.appendPlainText(f"→ Évaluation du modèle sur le dossier test ({test_folder_path})...")
        self.pbar.setValue(0)
        self._clear_evaluation_display() # Clear previous results before new evaluation
        
        opts = (self.chk_c.isChecked(), self.chk_d.isChecked(),
                self.chk_v.isChecked(), self.chk_m.isChecked())
                
        # Utilisation du TestEvaluationThread défini ci-dessous
        self.test_evaluator = TestEvaluationThread(
            test_folder_path, 
            self.models[self.curr_model_name], 
            self.trained_labels, # Pass labels from training for consistency
            opts, 
            self.threshold,
            self
        )
        self.test_evaluator.progress.connect(self.pbar.setValue)
        self.test_evaluator.finished.connect(self._test_evaluation_ready)
        self.test_evaluator.start()

    @pyqtSlot(object)
    def _test_evaluation_ready(self, evaluation_results):
        self.pbar.setValue(100)
        if "error" in evaluation_results:
            QMessageBox.critical(self, "Erreur d'évaluation", evaluation_results["error"])
            self._clear_evaluation_display()
            return
            
        # Store results for the current model
        self.g_rep[self.curr_model_name] = evaluation_results 
        
        # Update GUI with evaluation results
        self._populate_tables()
        self._update_stats_overview()
        self._plot_det_curve()
        self._plot_tfa_curve()
        self._plot_tfr_curve()
 
       

        
        # Enable export actions
        for act in [self.act_export_csv, self.act_export_pdf, self.act_export_ppt]:
            act.setEnabled(True)
            
        self.txt_logs.appendPlainText("✓ Évaluation sur le dossier test terminée.")
        if PRETTYTABLE_OK:
            t = PrettyTable()
            t.field_names = ["Métrique (Test)", "Valeur"]
            if "identification_rate" in evaluation_results:
                 t.add_row(["Taux d'identification", f"{evaluation_results['identification_rate']:.4f}"])
            # Utilisation des clés corrigées
            if "tfa_avg" in evaluation_results:
                 t.add_row(["TFA moyen", f"{evaluation_results['tfa_avg']:.4f}"])
            if "tfr_avg" in evaluation_results:
                 t.add_row(["TFR moyen", f"{evaluation_results['tfr_avg']:.4f}"])
            if "det" in evaluation_results and evaluation_results["det"]:
                 t.add_row(["TEE", f"{evaluation_results['det']['tee']:.4f}"])
            self.txt_logs.appendPlainText(t.get_string())
        else:
             if "identification_rate" in evaluation_results:
                 self.txt_logs.appendPlainText(f"  Taux d'identification: {evaluation_results['identification_rate']:.4f}")
             # Utilisation des clés corrigées
             if "tfa_avg" in evaluation_results:
                 self.txt_logs.appendPlainText(f"  TFA moyen: {evaluation_results['tfa_avg']:.4f}")
             if "tfr_avg" in evaluation_results:
                 self.txt_logs.appendPlainText(f"  TFR moyen: {evaluation_results['tfr_avg']:.4f}")
             if "det" in evaluation_results and evaluation_results["det"]:
                 self.txt_logs.appendPlainText(f"  TEE: {evaluation_results['det']['tee']:.4f}")

    def _clear_evaluation_display(self):
        """Clears all tables, plots, and stats related to evaluation."""
        self.g_rep = {} # Clear stored results
        # Clear tables
        self.tbl_metrics.setRowCount(0); self.tbl_metrics.setColumnCount(0)
        self.tbl_cm.setRowCount(0); self.tbl_cm.setColumnCount(0)
        # Clear stats text
        self.txt_stats.clear()
        # Clear plots
        for canvas in [self.canvas_tfa, self.canvas_tfr, self.canvas_det]:
            if hasattr(canvas, 'figure'):
                try:
                    fig = canvas.figure
                    fig.clear()
                    ax = fig.add_subplot(111)
                    ax.text(0.5, 0.5, "Évaluation requise", 
                           horizontalalignment='center',
                           verticalalignment='center',
                           transform=ax.transAxes,
                           fontsize=12)
                    ax.set_xticks([])
                    ax.set_yticks([])
                    canvas.draw()
                except Exception as e:
                    print(f"Erreur lors du nettoyage du canvas: {e}") # Debug

    def _populate_tables(self):
        """Populates the metrics and confusion matrix tables."""
        if not self.g_rep or self.curr_model_name not in self.g_rep: return
        rep = self.g_rep[self.curr_model_name]
        if not rep or 'rep_dict' not in rep or 'cm' not in rep or 'labels' not in rep: return

        rep_dict = rep['rep_dict']
        cm = rep['cm']
        labels = rep['labels']

        # --- Metrics Table ---
        metrics = ['precision', 'recall', 'f1-score', 'support']
        self.tbl_metrics.setRowCount(len(labels))
        self.tbl_metrics.setColumnCount(len(metrics))
        self.tbl_metrics.setHorizontalHeaderLabels(metrics)
        self.tbl_metrics.setVerticalHeaderLabels(labels)
        for i, label in enumerate(labels):
            if label in rep_dict:
                for j, metric in enumerate(metrics):
                    val = rep_dict[label][metric]
                    item = QTableWidgetItem(f"{val:.4f}" if isinstance(val, float) else str(val))
                    item.setTextAlignment(Qt.AlignCenter)
                    self.tbl_metrics.setItem(i, j, item)
        self.tbl_metrics.resizeColumnsToContents()
        self.tbl_metrics.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)

        # --- Confusion Matrix Table ---
        self.tbl_cm.setRowCount(len(labels))
        self.tbl_cm.setColumnCount(len(labels))
        self.tbl_cm.setHorizontalHeaderLabels(labels)
        self.tbl_cm.setVerticalHeaderLabels(labels)
        for i in range(len(labels)):
            for j in range(len(labels)):
                item = QTableWidgetItem(str(cm[i, j]))
                item.setTextAlignment(Qt.AlignCenter)
                self.tbl_cm.setItem(i, j, item)
        self.tbl_cm.resizeColumnsToContents()
        self.tbl_cm.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
        self.tbl_cm.verticalHeader().setSectionResizeMode(QHeaderView.Stretch)

    def _update_stats_overview(self):
        """Met à jour l'affichage des statistiques globales."""
        if not self.g_rep or self.curr_model_name not in self.g_rep:
            self.txt_stats.clear()
            return
            
        rep = self.g_rep[self.curr_model_name]
        stats_text = "Statistiques globales sur l'ensemble de test :\n\n"
        
        # Affichage des métriques avec nomenclature cohérente
        if "identification_rate" in rep:
            stats_text += f"Taux d'identification : {rep['identification_rate']:.4f}\n"
        
        # Utilisation des clés corrigées
        if "tfa_avg" in rep:
            stats_text += f"TFA moyen (Taux de Fausse Acceptation) : {rep['tfa_avg']:.4f}\n"
        
        if "tfr_avg" in rep:
            stats_text += f"TFR moyen (Taux de Faux Rejet) : {rep['tfr_avg']:.4f}\n"
        
        if "det" in rep and rep["det"]:
            det_data = rep["det"]
            stats_text += f"TEE (Taux d'Égale Erreur) : {det_data['tee']:.4f}\n"
            stats_text += f"Seuil optimal au TEE : {det_data['seuils'][det_data['tee_idx']]:.2f}\n"
        
        self.txt_stats.setPlainText(stats_text)

    def _plot_roc(self):
        """Affiche la courbe ROC à partir des données d'évaluation."""
        if not self.g_rep or self.curr_model_name not in self.g_rep: 
            self._display_no_data_message(self.canvas_roc, "Pas de données ROC disponibles")
            return
            
        rep = self.g_rep[self.curr_model_name]
        if "roc" not in rep or not rep["roc"]: 
            self._display_no_data_message(self.canvas_roc, "Pas de données ROC disponibles")
            return
            
        tfa, tpr, roc_auc = rep["roc"] # Utilise fpr comme TFA
        
        fig = self.canvas_roc.figure
        fig.clear()
        ax = fig.add_subplot(111)
        
        # Tracé de la courbe ROC avec une meilleure visibilité
        ax.plot(tfa, tpr, 'b-', linewidth=2, label=f'AUC = {roc_auc:.4f}')
        ax.plot([0, 1], [0, 1], 'k--', alpha=0.5)  # Ligne de référence
        
        # Amélioration de l'apparence et des informations
        ax.set_title('Courbe ROC')
        ax.set_xlabel('Taux de Fausse Acceptation (TFA)')
        ax.set_ylabel('Taux de Vraie Acceptation (1 - TFR)')
        ax.legend(loc='lower right')
        ax.grid(True, alpha=0.3)
        
        # Mise à jour du graphique
        self.canvas_roc.draw()
    def _plot_tfa_curve(self):
     """Trace TFA = FAR en fonction du seuil."""
     if not self._det_ready(): return
     det = self.g_rep[self.curr_model_name]["det"]
     seuils, tfa = det["seuils"], det["tfa"]      # déjà calculés
     fig = self.canvas_tfa.figure; fig.clear()
     ax  = fig.add_subplot(111)
     ax.plot(seuils, tfa, lw=2, color="tab:blue", label="TFA") 
     ax.set_xlabel("Seuil"); ax.set_ylabel("TFA")
     ax.set_title("Courbe TFA"); ax.grid(True)
     self.canvas_tfa.draw()

    def _plot_tfr_curve(self):
     """Trace TFR = FRR en fonction du seuil."""
     if not self._det_ready(): return
     det = self.g_rep[self.curr_model_name]["det"]
     seuils, tfr = det["seuils"], det["tfr"]
     fig = self.canvas_tfr.figure; fig.clear()
     ax  = fig.add_subplot(111)
     
     ax.plot(seuils, tfr, lw=2, color="tab:red", label="TFR")
     ax.set_xlabel("Seuil"); ax.set_ylabel("TFR")
     ax.set_title("Courbe TFR"); ax.grid(True)
     self.canvas_tfr.draw()

    def _det_ready(self):
     """Petit helper interne pour vérifier la présence des données DET."""
     if not self.g_rep or self.curr_model_name not in self.g_rep:
        self._display_no_data_message(self.canvas_tfa, "Pas de données DET")
        self._display_no_data_message(self.canvas_tfr, "Pas de données DET")
        return False
     if "det" not in self.g_rep[self.curr_model_name]:
        return False
     return True
    def _plot_det_curve(self):
        """Affiche la courbe DET (TFA/TFR/TEE) à partir des données d'évaluation."""
        if not self.g_rep or self.curr_model_name not in self.g_rep:
            self._display_no_data_message(self.canvas_det, "Pas de données DET disponibles")
            return
            
        rep = self.g_rep[self.curr_model_name]
        if "det" not in rep or not rep["det"]:
            self._display_no_data_message(self.canvas_det, "Pas de données DET disponibles")
            return
            
        det_data = rep["det"]
        
        # Récupération des données avec nomenclature cohérente
        tfa = det_data["tfa"] if "tfa" in det_data else det_data["fpr"]  # Taux de Fausse Acceptation
        tfr = det_data["tfr"] if "tfr" in det_data else det_data["frr"]  # Taux de Faux Rejet
        tee = det_data["tee"] if "tee" in det_data else det_data["eer"]  # Taux d'Égale Erreur
        tee_idx = det_data["tee_idx"] if "tee_idx" in det_data else det_data["eer_idx"]
        seuils = det_data["seuils"] if "seuils" in det_data else det_data["x_scale"]

        
        
        # Préparation du graphique
        fig = self.canvas_det.figure
        fig.clear()
        ax = fig.add_subplot(111)
        
        # Tracé des courbes avec une meilleure visibilité
        ax.plot(seuils, tfa, 'b-', linewidth=2, label='TFA (Taux de Fausse Acceptation)')
        ax.plot(seuils, tfr, 'r-', linewidth=2, label='TFR (Taux de Faux Rejet)')
        
        # Mise en évidence du point TEE
        ax.plot(seuils[tee_idx], tee, 'ko', markersize=8, label=f'TEE = {tee:.4f}')
        ax.axvline(x=seuils[tee_idx], color='k', linestyle='--', alpha=0.5)
        ax.axhline(y=tee, color='k', linestyle='--', alpha=0.5)
        
        # Amélioration de l'apparence et des informations
        ax.set_title('Courbe DET : TFA et TFR en fonction du seuil')
        ax.set_xlabel('Seuil de décision')
        ax.set_ylabel('Taux d\'erreur')
        ax.legend(loc='best')
        ax.grid(True, alpha=0.3)
        
        # Ajout d'annotations pour le point TEE
        ax.annotate(f'TEE = {tee:.4f}',
                    xy=(seuils[tee_idx], tee),
                    xytext=(seuils[tee_idx] + 5, tee + 0.05),
                    arrowprops=dict(facecolor='black', shrink=0.05, width=1.5),
                    fontsize=9)
        
        # Mise à jour du graphique
        self.canvas_det.draw()
        
    def _display_no_data_message(self, canvas, message):
        """Affiche un message d'erreur sur le canvas."""
        fig = canvas.figure
        fig.clear()
        ax = fig.add_subplot(111)
        ax.text(0.5, 0.5, message, 
               horizontalalignment='center',
               verticalalignment='center',
               transform=ax.transAxes,
               fontsize=12,
               color='red')
        ax.set_xticks([])
        ax.set_yticks([])
        canvas.draw()

    def _update_lda_overview(self):
        """Updates the LDA model overview tab (e.g., coefficients)."""
        if not self.models or self.curr_model_name != "LDA":
            # Clear if no model or not LDA
            self.tbl_feats.setRowCount(0)
            self.tbl_feats.setColumnCount(0)
            return

        try:
            lda_model = self.models["LDA"].named_steps['customlda']
            if hasattr(lda_model, 'scalings_'):
                scalings = lda_model.scalings_
                # Assuming feature names are not stored, use indices
                num_components = scalings.shape[1]
                num_features = scalings.shape[0]
                
                self.tbl_feats.setRowCount(num_features)
                self.tbl_feats.setColumnCount(num_components)
                self.tbl_feats.setHorizontalHeaderLabels([f"LD{i+1}" for i in range(num_components)])
                self.tbl_feats.setVerticalHeaderLabels([f"Feat_{i}" for i in range(num_features)])
                
                for i in range(num_features):
                    for j in range(num_components):
                        item = QTableWidgetItem(f"{scalings[i, j]:.4f}")
                        item.setTextAlignment(Qt.AlignCenter)
                        self.tbl_feats.setItem(i, j, item)
                self.tbl_feats.resizeColumnsToContents()
            else:
                 self.tbl_feats.setRowCount(0)
                 self.tbl_feats.setColumnCount(0)
        except Exception as e:
            print(f"Erreur lors de la mise à jour de l'aperçu LDA: {e}")
            self.tbl_feats.setRowCount(0)
            self.tbl_feats.setColumnCount(0)

    def _prompt_open_image(self):
        fpath, _ = QFileDialog.getOpenFileName(self, "Ouvrir image", "", "Images (*.png *.jpg *.bmp *.tif)")
        if fpath:
            self._load_image(fpath)

    def _load_image(self, fpath):
        self.current_image_path = fpath
        
        # Charger l'image en couleur (BGR par défaut pour OpenCV)
        img_color = cv2.imread(fpath) 
        
        if img_color is None: 
            QMessageBox.critical(self,"Erreur","Impossible de charger l'image."); return
        
        # Afficher l'image originale en couleur dans self.lbl_orig
        # La fonction numpy_to_qpixmap doit pouvoir gérer les images couleur (BGR, 3 canaux)
        # La version dans pipeline.py semble le faire.
        self.lbl_orig.setPixmap(numpy_to_qpixmap(cv2.resize(img_color, (400,320)))) # Taille modifiée
        
        # Convertir l'image en niveaux de gris pour le pipeline de traitement
        # car les fonctions comme binarize, skeletonize, lbp attendent du niveaux de gris.
        img_gray_for_pipeline = cv2.cvtColor(img_color, cv2.COLOR_BGR2GRAY)
        
        opts = (self.chk_c.isChecked(), self.chk_d.isChecked(),
                self.chk_v.isChecked(), self.chk_m.isChecked())
        
        # Utilisation de pipeline_full avec l'image en NIVEAUX DE GRIS
        # Assurez-vous que pipeline_full est importé et disponible
        try:
            # b, s, lbp_img, _ = pipeline_full(img_gray_for_pipeline, *opts)
            # Si pipeline_full est une méthode de la classe, utilisez self.pipeline_full
            # S'il vient d'une importation, utilisez directement pipeline_full
            # Pour cet exemple, je suppose qu'il est importé globalement ou accessible.
            # Si ce n'est pas le cas, ajustez l'appel.
            
            # Appel à la fonction pipeline_full (qui doit être accessible ici)
            # Par exemple, si elle est dans pipeline.py et importée :
            # from pipeline import pipeline_full # au début du fichier gui.py
            b, s, lbp_img, _ = pipeline_full(img_gray_for_pipeline, *opts)


        except NameError as e:
            QMessageBox.critical(self, "Erreur de Pipeline", f"La fonction 'pipeline_full' n'est pas définie ou importée correctement.\n{e}")
            return
        except Exception as e:
            QMessageBox.critical(self, "Erreur de Pipeline", f"Erreur lors de l'exécution de pipeline_full: {e}")
            return

        # Afficher les images traitées (qui sont en niveaux de gris/binaires)
        self.lbl_bin.setPixmap(numpy_to_qpixmap(cv2.resize(b,(200,200))))
        self.lbl_skel.setPixmap(numpy_to_qpixmap(cv2.resize(s,(200,200))))
        self.lbl_lbp.setPixmap(numpy_to_qpixmap(cv2.resize(lbp_img,(200,200))))
        
        # Activer l'analyse si un modèle est entraîné
        if self.models:
            self.act_analyse.setEnabled(True)

    def _analyse(self):
    # 1) Vérification du chemin d'image
     if not hasattr(self, 'current_image_path') or not self.current_image_path:
        QMessageBox.warning(self, "Erreur", "Veuillez d'abord charger une image.")
        return

    # 2) Lecture de l'image
     img = cv2.imread(self.current_image_path)
     if img is None:
        QMessageBox.warning(self, "Erreur", "Impossible de lire l'image sélectionnée.")
        return

    # 3) Validation "retina-like"
     if not self._validate_retina_image_advanced(img):
        QMessageBox.warning(
            self,
            "Image invalide",
            "L'image sélectionnée ne semble pas être une image de rétine valide."
        )
        return  # on n'analyse pas si ce n'est pas une rétine

    # 4) Vérifier qu'on a bien un modèle entraîné
     if not self.models:
        QMessageBox.warning(self, "Erreur", "Veuillez d'abord entraîner un modèle.")
        return

    # 5) Choix du mode et gestion du claim en mode auth
     mode = "ident" if self.rb_ident.isChecked() else "auth"
     claim = self.line_claim.text().strip()
     if mode == "auth" and not claim:
        QMessageBox.warning(self, "Erreur", "Veuillez spécifier l'identité revendiquée.")
        return

    # 6) Configuration de la barre de progression et des options
     self.pbar.setValue(0)
     opts = (
        self.chk_c.isChecked(),
        self.chk_d.isChecked(),
        self.chk_v.isChecked(),
        self.chk_m.isChecked()
    )

    # 7) Lancement du thread d'analyse
     current_eval_results = self.g_rep.get(self.curr_model_name, {})
     self.analyser = AnalyseThread(
        self.current_image_path,
        self.models[self.curr_model_name],
        opts,
        mode,
        claim,
        current_eval_results,
        self.threshold,
        self
    )
     self.analyser.progress.connect(self.pbar.setValue)
     self.analyser.done.connect(self._analyse_done)
     self.analyser.start()


    @pyqtSlot(str)
    def _analyse_done(self, msg):
        self.pbar.setValue(100)
        self.txt_logs.appendPlainText(f"→ Analyse: {msg}")
       

    def _switch_method(self, name):
        self.curr_model_name = name
        # Re-display results for the selected model if they exist
        if self.g_rep and name in self.g_rep:
            self._populate_tables()
            self._update_stats_overview()
            self._plot_tfa_curve()
            self._plot_tfr_curve()
            self._plot_det_curve()
            self._update_lda_overview()
        else:
            self._clear_evaluation_display()

    def _export_csv(self):
        if not self.g_rep or self.curr_model_name not in self.g_rep: return
        fpath, _ = QFileDialog.getSaveFileName(self, "Exporter CSV", "", "CSV (*.csv)")
        if not fpath: return
        try:
            rep = self.g_rep[self.curr_model_name]
            df = pd.DataFrame(rep['rep_dict']).transpose()
            df.to_csv(fpath)
            self.txt_logs.appendPlainText(f"✓ Export CSV: {fpath}")
        except Exception as e: QMessageBox.critical(self,"Erreur Export",str(e))

    def _export_pdf(self):
        if not PDF_OK: QMessageBox.critical(self,"Erreur","Bibliothèque ReportLab non trouvée."); return
        if not self.g_rep or self.curr_model_name not in self.g_rep: return
        fpath, _ = QFileDialog.getSaveFileName(self, "Exporter PDF", "", "PDF (*.pdf)")
        if not fpath: return
        try:
            c = pdf_canvas.Canvas(fpath, pagesize=letter)
            w, h = letter
            rep = self.g_rep[self.curr_model_name]
            
            # Titre
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, h - 50, f"Rapport d'évaluation - Modèle {self.curr_model_name}")
            
            # Statistiques globales
            c.setFont("Helvetica-Bold", 12)
            c.drawString(50, h - 80, "Statistiques Globales (Test)")
            c.setFont("Helvetica", 10)
            y_pos = h - 100
            if "identification_rate" in rep: c.drawString(60, y_pos, f"Taux d'identification: {rep['identification_rate']:.4f}"); y_pos -= 15
            if "tfa_avg" in rep: c.drawString(60, y_pos, f"TFA moyen: {rep['tfa_avg']:.4f}"); y_pos -= 15
            if "tfr_avg" in rep: c.drawString(60, y_pos, f"TFR moyen: {rep['tfr_avg']:.4f}"); y_pos -= 15
            if "det" in rep and rep["det"]: c.drawString(60, y_pos, f"TEE: {rep['det']['tee']:.4f}"); y_pos -= 15
            
            # Ajouter les graphiques (enregistrés temporairement)
            img_roc_path = self._save_canvas_to_temp(self.canvas_roc)
            img_det_path = self._save_canvas_to_temp(self.canvas_det)
            
            if img_roc_path:
                c.drawImage(img_roc_path, 50, y_pos - 220, width=250, height=200)
                os.remove(img_roc_path)
            if img_det_path:
                c.drawImage(img_det_path, 310, y_pos - 220, width=250, height=200)
                os.remove(img_det_path)

            # Ajouter les tableaux (simplifié)
            # ... (l'ajout de tableaux complexes nécessite plus de code avec reportlab.platypus)
            
            c.save()
            self.txt_logs.appendPlainText(f"✓ Export PDF: {fpath}")
        except Exception as e: QMessageBox.critical(self,"Erreur Export",str(e))

    def _export_ppt(self):
        if not PPTX_OK: QMessageBox.critical(self,"Erreur","Bibliothèque python-pptx non trouvée."); return
        if not self.g_rep or self.curr_model_name not in self.g_rep: return
        fpath, _ = QFileDialog.getSaveFileName(self, "Exporter PPTX", "", "PPTX (*.pptx)")
        if not fpath: return
        try:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5] # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            rep = self.g_rep[self.curr_model_name]

            # Titre
            title = slide.shapes.title
            title.text = f"Rapport d'évaluation - Modèle {self.curr_model_name}"

            # Ajouter les graphiques
            img_roc_path = self._save_canvas_to_temp(self.canvas_roc)
            img_det_path = self._save_canvas_to_temp(self.canvas_det)
            
            if img_roc_path:
                slide.shapes.add_picture(img_roc_path, Inches(0.5), Inches(1.5), width=Inches(4.5))
                os.remove(img_roc_path)
            if img_det_path:
                slide.shapes.add_picture(img_det_path, Inches(5.0), Inches(1.5), width=Inches(4.5))
                os.remove(img_det_path)

            # Ajouter les statistiques (simplifié)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9.0), Inches(1.5))
            tf = txBox.text_frame
            tf.text = "Statistiques Globales (Test):\n"
            if "identification_rate" in rep: tf.add_paragraph(f"Taux d'identification: {rep['identification_rate']:.4f}")
            if "tfa_avg" in rep: tf.add_paragraph(f"TFA moyen: {rep['tfa_avg']:.4f}")
            if "tfr_avg" in rep: tf.add_paragraph(f"TFR moyen: {rep['tfr_avg']:.4f}")
            if "det" in rep and rep["det"]: tf.add_paragraph(f"TEE: {rep['det']['tee']:.4f}")

            prs.save(fpath)
            self.txt_logs.appendPlainText(f"✓ Export PPTX: {fpath}")
        except Exception as e: QMessageBox.critical(self,"Erreur Export",str(e))

    def _save_canvas_to_temp(self, canvas):
        """Saves a matplotlib canvas to a temporary PNG file."""
        try:
            img_data = BytesIO()
            canvas.figure.savefig(img_data, format='png', bbox_inches='tight')
            img_data.seek(0)
            temp_path = "temp_plot.png"
            with open(temp_path, 'wb') as f:
                f.write(img_data.read())
            return temp_path
        except Exception as e:
            print(f"Erreur sauvegarde canvas: {e}")
            return None

    def _stylesheet(self):
        return """
            QMainWindow { background-color: #f0f0f0; }
            QGroupBox { border: 1px solid #cccccc; border-radius: 5px; margin-top: 10px; }
            QGroupBox::title { subcontrol-origin: margin; subcontrol-position: top center; padding: 0 3px; background-color: #f0f0f0; }
            QLabel { font-size: 11pt; }
            QPushButton { background-color: #4CAF50; color: white; padding: 8px; border-radius: 4px; font-size: 10pt; }
            QPushButton:hover { background-color: #45a049; }
            QPushButton:disabled { background-color: #cccccc; }
            QComboBox { padding: 5px; border: 1px solid #ccc; border-radius: 3px; }
            QLineEdit { padding: 5px; border: 1px solid #ccc; border-radius: 3px; }
            QPlainTextEdit { border: 1px solid #ccc; border-radius: 3px; font-family: Consolas, monospace; }
            QProgressBar { border: 1px solid #ccc; border-radius: 3px; text-align: center; }
            QProgressBar::chunk { background-color: #4CAF50; }
            QTableWidget { border: 1px solid #ccc; gridline-color: #e0e0e0; }
            QHeaderView::section { background-color: #e8e8e8; padding: 4px; border: 1px solid #ccc; }
            QTabWidget::pane { border: 1px solid #ccc; border-top: none; }
            QTabBar::tab { background: #e0e0e0; padding: 8px; border: 1px solid #ccc; border-bottom: none; border-top-left-radius: 4px; border-top-right-radius: 4px; }
            QTabBar::tab:selected { background: #f0f0f0; }
            QToolBar { background: #e0e0e0; border: none; padding: 5px; }
            QAction { color: #333; }
        """

# --- Thread pour l'évaluation sur le dossier Test --- 
# (Doit être défini ici car il utilise des éléments de la GUI)
class TestEvaluationThread(QThread):
    progress = pyqtSignal(int)
    finished = pyqtSignal(object) # Renvoie le dictionnaire de résultats

    def __init__(self, test_folder, model, trained_labels, opts, threshold, parent=None):
        super().__init__(parent)
        self.test_folder = test_folder
        self.model = model
        self.trained_labels = trained_labels
        self.opts = opts
        self.threshold = threshold

    def run(self):
        try:
            files = [
                f for f in sorted(os.listdir(self.test_folder))
                if f.lower().endswith(('.png', '.jpg', '.bmp', '.tif'))
            ]
            X_test, y_test = [], []
            
            for i, fname in enumerate(files):
                img = cv2.imread(os.path.join(self.test_folder, fname), cv2.IMREAD_GRAYSCALE)
                _, _, _, feats = pipeline_full(img, *self.opts) # Utilise pipeline_full de pipeline_corrige
                X_test.append(feats)
                label = fname.split('_')[1].split('.')[0].lower()
                y_test.append(label)
                self.progress.emit(int((i + 1) / len(files) * 50)) # Progression jusqu'à 50%

            X_test = np.array(X_test)
            y_test = np.array(y_test)

            # Évaluation classique (utilise evaluate de pipeline_corrige)
            results = evaluate(self.model, X_test, y_test, self.trained_labels, self.threshold)
            self.progress.emit(70)

            # Calcul Courbe ROC (utilise compute_roc de pipeline_corrige)
            roc_data = compute_roc(self.model, X_test, y_test, self.trained_labels)
            if roc_data:
                results["roc"] = roc_data
            self.progress.emit(85)

            # Calcul Courbe DET (TFA/TFR/TEE) (utilise compute_det_curve de pipeline_corrige)
            det_data = compute_det_curve(self.model, X_test, y_test, self.trained_labels)
            if det_data:
                results["det"] = det_data
                # Ajouter les moyennes pour les stats globales
                results["tfa_avg"] = det_data.get("tfa_avg", 0)
                results["tfr_avg"] = det_data.get("tfr_avg", 0)
            self.progress.emit(95)
            
            # Calcul du taux d'identification (si nécessaire)
            y_pred = self.model.predict(X_test)
            identification_rate = np.mean(y_pred == y_test)
            results["identification_rate"] = identification_rate

            self.progress.emit(100)
            self.finished.emit(results)

        except Exception as e:
            import traceback
            print(traceback.format_exc())
            self.finished.emit({"error": str(e)})

def main():
    import sys
    from PyQt5.QtWidgets import QApplication
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())

if __name__ == '__main__':
    main()

if __name__ == '__main__':
    app = QApplication(sys.argv)
    win = MainWindow()
    win.show()
    sys.exit(app.exec_())
